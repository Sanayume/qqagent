"""
文件读取器 - 支持多种文件格式的内容提取

支持的格式:
- PDF: 提取文本内容（需要 pymupdf 或 pdfplumber）
- 图片: 返回 base64 编码（供后续处理）
- 文本/代码: 直接读取内容
- Office (docx/xlsx/pptx): 提取文本内容
"""

import base64
from pathlib import Path
from dataclasses import dataclass
from enum import Enum


class FileType(Enum):
    """文件类型"""
    PDF = "pdf"
    IMAGE = "image"
    TEXT = "text"
    CODE = "code"
    DOCX = "docx"
    XLSX = "xlsx"
    PPTX = "pptx"
    UNKNOWN = "unknown"


@dataclass
class FileContent:
    """文件内容"""
    success: bool
    file_type: FileType
    content: str  # 文本内容或 base64
    is_base64: bool = False
    mime_type: str = ""  # MIME 类型（用于图片/PDF）
    error: str = ""

    def to_llm_format(self) -> str:
        """转换为 LLM 可读格式（文本工具返回值）"""
        if not self.success:
            return f"错误: {self.error}"

        if self.is_base64:
            # 对于 base64 内容，返回提示信息
            # 实际的多模态处理需要调用方单独处理
            return f"[{self.file_type.value.upper()} 文件已读取，base64 长度: {len(self.content)} 字符，MIME: {self.mime_type}]"

        return self.content


# 文件扩展名映射
EXT_TO_TYPE: dict[str, FileType] = {
    # PDF
    ".pdf": FileType.PDF,
    # 图片
    ".png": FileType.IMAGE,
    ".jpg": FileType.IMAGE,
    ".jpeg": FileType.IMAGE,
    ".gif": FileType.IMAGE,
    ".webp": FileType.IMAGE,
    ".bmp": FileType.IMAGE,
    # 文本
    ".txt": FileType.TEXT,
    ".md": FileType.TEXT,
    ".markdown": FileType.TEXT,
    ".rst": FileType.TEXT,
    ".log": FileType.TEXT,
    ".csv": FileType.TEXT,
    ".json": FileType.TEXT,
    ".yaml": FileType.TEXT,
    ".yml": FileType.TEXT,
    ".xml": FileType.TEXT,
    ".html": FileType.TEXT,
    ".htm": FileType.TEXT,
    ".ini": FileType.TEXT,
    ".cfg": FileType.TEXT,
    ".conf": FileType.TEXT,
    ".toml": FileType.TEXT,
    # 代码
    ".py": FileType.CODE,
    ".js": FileType.CODE,
    ".ts": FileType.CODE,
    ".jsx": FileType.CODE,
    ".tsx": FileType.CODE,
    ".java": FileType.CODE,
    ".c": FileType.CODE,
    ".cpp": FileType.CODE,
    ".h": FileType.CODE,
    ".hpp": FileType.CODE,
    ".cs": FileType.CODE,
    ".go": FileType.CODE,
    ".rs": FileType.CODE,
    ".rb": FileType.CODE,
    ".php": FileType.CODE,
    ".swift": FileType.CODE,
    ".kt": FileType.CODE,
    ".scala": FileType.CODE,
    ".r": FileType.CODE,
    ".sql": FileType.CODE,
    ".sh": FileType.CODE,
    ".bash": FileType.CODE,
    ".zsh": FileType.CODE,
    ".ps1": FileType.CODE,
    ".bat": FileType.CODE,
    ".cmd": FileType.CODE,
    ".vue": FileType.CODE,
    ".svelte": FileType.CODE,
    # Office
    ".docx": FileType.DOCX,
    ".xlsx": FileType.XLSX,
    ".pptx": FileType.PPTX,
}

# 最大文本文件大小 (1MB)
MAX_TEXT_SIZE = 1 * 1024 * 1024

# 最大二进制文件大小 (10MB) - 用于 PDF、图片
MAX_BINARY_SIZE = 10 * 1024 * 1024

# 最大输出内容长度 (50000 字符，约 50KB 文本)
MAX_CONTENT_LENGTH = 50000


def _truncate_content(content: str, max_length: int = MAX_CONTENT_LENGTH) -> str:
    """截断过长的内容"""
    if len(content) <= max_length:
        return content
    return content[:max_length] + f"\n\n... [内容过长，已截断。总长度: {len(content)} 字符，显示前 {max_length} 字符]"


def get_file_type(file_path: Path) -> FileType:
    """根据扩展名判断文件类型"""
    ext = file_path.suffix.lower()
    return EXT_TO_TYPE.get(ext, FileType.UNKNOWN)


def read_file(file_path: str) -> FileContent:
    """读取文件内容

    Args:
        file_path: 文件路径

    Returns:
        FileContent 对象
    """
    path = Path(file_path)

    if not path.exists():
        return FileContent(
            success=False,
            file_type=FileType.UNKNOWN,
            content="",
            error=f"文件不存在: {file_path}"
        )

    file_type = get_file_type(path)

    try:
        if file_type == FileType.PDF:
            return _read_pdf(path)
        elif file_type == FileType.IMAGE:
            return _read_image(path)
        elif file_type in (FileType.TEXT, FileType.CODE):
            return _read_text(path, file_type)
        elif file_type == FileType.DOCX:
            return _read_docx(path)
        elif file_type == FileType.XLSX:
            return _read_xlsx(path)
        elif file_type == FileType.PPTX:
            return _read_pptx(path)
        else:
            return FileContent(
                success=False,
                file_type=FileType.UNKNOWN,
                content="",
                error=f"不支持的文件格式: {path.suffix}"
            )
    except Exception as e:
        return FileContent(
            success=False,
            file_type=file_type,
            content="",
            error=f"读取失败: {e}"
        )


def _read_pdf(path: Path) -> FileContent:
    """读取 PDF 文件，尝试提取文本内容"""
    # 检查文件大小
    size = path.stat().st_size
    if size > MAX_BINARY_SIZE:
        return FileContent(
            success=False,
            file_type=FileType.PDF,
            content="",
            error=f"文件太大: {size / 1024 / 1024:.1f}MB > 10MB 限制"
        )

    # 尝试使用 PyMuPDF 提取文本
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(path)
        text_parts = []
        for page_num, page in enumerate(doc, 1):
            text = page.get_text().strip()
            if text:
                text_parts.append(f"=== 第 {page_num} 页 ===\n{text}")
        doc.close()

        if text_parts:
            content = "\n\n".join(text_parts)
            return FileContent(
                success=True,
                file_type=FileType.PDF,
                content=_truncate_content(content)
            )
        else:
            return FileContent(
                success=True,
                file_type=FileType.PDF,
                content="(PDF 无可提取的文本内容，可能是扫描件或纯图片)"
            )
    except ImportError:
        pass  # PyMuPDF 未安装，尝试其他方法

    # 尝试使用 pdfplumber
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    text_parts.append(f"=== 第 {page_num} 页 ===\n{text.strip()}")

        if text_parts:
            content = "\n\n".join(text_parts)
            return FileContent(
                success=True,
                file_type=FileType.PDF,
                content=_truncate_content(content)
            )
        else:
            return FileContent(
                success=True,
                file_type=FileType.PDF,
                content="(PDF 无可提取的文本内容，可能是扫描件或纯图片)"
            )
    except ImportError:
        pass  # pdfplumber 未安装

    # 都没安装，返回提示
    return FileContent(
        success=False,
        file_type=FileType.PDF,
        content="",
        error="需要安装 PDF 解析库: pip install pymupdf 或 pip install pdfplumber"
    )


# 图片扩展名到 MIME 类型映射
EXT_TO_MIME: dict[str, str] = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
    ".bmp": "image/bmp",
}


def _read_image(path: Path) -> FileContent:
    """读取图片文件，返回 base64 和 MIME 类型"""
    # 检查文件大小
    size = path.stat().st_size
    if size > MAX_BINARY_SIZE:
        return FileContent(
            success=False,
            file_type=FileType.IMAGE,
            content="",
            error=f"图片太大: {size / 1024 / 1024:.1f}MB > 10MB 限制"
        )

    with open(path, "rb") as f:
        data = f.read()

    ext = path.suffix.lower()
    mime_type = EXT_TO_MIME.get(ext, "image/png")

    b64 = base64.b64encode(data).decode("utf-8")

    return FileContent(
        success=True,
        file_type=FileType.IMAGE,
        content=b64,
        is_base64=True,
        mime_type=mime_type
    )


def _read_text(path: Path, file_type: FileType) -> FileContent:
    """读取文本文件"""
    size = path.stat().st_size
    if size > MAX_TEXT_SIZE:
        return FileContent(
            success=False,
            file_type=file_type,
            content="",
            error=f"文件太大: {size / 1024 / 1024:.1f}MB > 1MB 限制"
        )

    # 尝试多种编码
    for encoding in ["utf-8", "gbk", "gb2312", "latin-1"]:
        try:
            with open(path, "r", encoding=encoding) as f:
                content = f.read()
            return FileContent(
                success=True,
                file_type=file_type,
                content=_truncate_content(content)
            )
        except UnicodeDecodeError:
            continue

    return FileContent(
        success=False,
        file_type=file_type,
        content="",
        error="无法解码文件内容"
    )


def _read_docx(path: Path) -> FileContent:
    """读取 Word 文档"""
    # 检查文件大小
    size = path.stat().st_size
    if size > MAX_BINARY_SIZE:
        return FileContent(
            success=False,
            file_type=FileType.DOCX,
            content="",
            error=f"文件太大: {size / 1024 / 1024:.1f}MB > 10MB 限制"
        )

    try:
        from docx import Document
    except ImportError:
        return FileContent(
            success=False,
            file_type=FileType.DOCX,
            content="",
            error="需要安装 python-docx: pip install python-docx"
        )

    doc = Document(path)
    result = []

    # 提取段落
    for p in doc.paragraphs:
        if p.text.strip():
            result.append(p.text)

    # 提取表格
    for table in doc.tables:
        table_rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            table_rows.append("\t".join(cells))
        if table_rows:
            result.append("--- 表格 ---")
            result.extend(table_rows)

    content = "\n".join(result)

    return FileContent(
        success=True,
        file_type=FileType.DOCX,
        content=_truncate_content(content) if content else "(文档为空)"
    )


def _read_xlsx(path: Path) -> FileContent:
    """读取 Excel 文件"""
    # 检查文件大小
    size = path.stat().st_size
    if size > MAX_BINARY_SIZE:
        return FileContent(
            success=False,
            file_type=FileType.XLSX,
            content="",
            error=f"文件太大: {size / 1024 / 1024:.1f}MB > 10MB 限制"
        )

    try:
        from openpyxl import load_workbook
    except ImportError:
        return FileContent(
            success=False,
            file_type=FileType.XLSX,
            content="",
            error="需要安装 openpyxl: pip install openpyxl"
        )

    wb = load_workbook(path, read_only=True, data_only=True)
    result = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        rows = []
        for row in sheet.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                row_str = "\t".join(str(cell) if cell is not None else "" for cell in row)
                rows.append(row_str)

        if rows:
            result.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows))

    wb.close()
    content = "\n\n".join(result)

    return FileContent(
        success=True,
        file_type=FileType.XLSX,
        content=_truncate_content(content) if content else "(表格为空)"
    )


def _read_pptx(path: Path) -> FileContent:
    """读取 PowerPoint 文件"""
    # 检查文件大小
    size = path.stat().st_size
    if size > MAX_BINARY_SIZE:
        return FileContent(
            success=False,
            file_type=FileType.PPTX,
            content="",
            error=f"文件太大: {size / 1024 / 1024:.1f}MB > 10MB 限制"
        )

    try:
        from pptx import Presentation
    except ImportError:
        return FileContent(
            success=False,
            file_type=FileType.PPTX,
            content="",
            error="需要安装 python-pptx: pip install python-pptx"
        )

    prs = Presentation(path)
    result = []

    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

        if texts:
            result.append(f"=== Slide {i} ===\n" + "\n".join(texts))

    content = "\n\n".join(result)

    return FileContent(
        success=True,
        file_type=FileType.PPTX,
        content=_truncate_content(content) if content else "(演示文稿为空)"
    )


def get_supported_extensions() -> list[str]:
    """获取支持的文件扩展名列表"""
    return list(EXT_TO_TYPE.keys())
